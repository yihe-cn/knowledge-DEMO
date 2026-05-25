"""文档 loader。支持 pptx 与文本型 pdf；docx 接口预留。"""
from __future__ import annotations

import re
from dataclasses import dataclass, field
from pathlib import Path


@dataclass
class RawSection:
    text: str
    meta: dict = field(default_factory=dict)  # 例如 {"slide_index": 3, "title": "..."}


def load_pptx(path: Path) -> list[RawSection]:
    from pptx import Presentation  # type: ignore

    prs = Presentation(str(path))
    sections: list[RawSection] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = ""
        body_parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                line = "".join(run.text for run in para.runs).strip()
                if not line:
                    continue
                # 第一段非空且短，视为标题
                if not title and len(line) <= 40:
                    title = line
                else:
                    body_parts.append(line)
        # speaker notes
        notes_text = ""
        try:
            if slide.has_notes_slide:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
        except Exception:
            notes_text = ""
        composed = "\n".join([title] + body_parts + ([f"[备注]{notes_text}"] if notes_text else [])).strip()
        if composed:
            sections.append(RawSection(text=composed, meta={"slide_index": idx, "title": title}))
    return sections


_BOILERPLATE_RE = re.compile(
    r"©|copyright|confidential|all right[s]? reserved|页\s*\d+\s*/|^\d+\s*$",
    re.IGNORECASE,
)


def _make_title(text: str) -> str:
    for line in text.splitlines():
        line = line.strip()
        if not line:
            continue
        # 跳过版权/页眉页脚噪声行，避免 title 全是 "© Copyright ..."
        if _BOILERPLATE_RE.search(line):
            continue
        return line[:40]
    return ""


def load_pdf(path: Path) -> list[RawSection]:
    from pypdf import PdfReader  # type: ignore

    reader = PdfReader(str(path))
    sections: list[RawSection] = []
    for idx, page in enumerate(reader.pages, start=1):
        text = (page.extract_text() or "").strip()
        if not text:
            continue
        sections.append(
            RawSection(text=text, meta={"page_index": idx, "title": _make_title(text)})
        )
    if not sections:
        raise ValueError("PDF contains no extractable text")
    return sections


def load_text(path: Path) -> list[RawSection]:
    text = path.read_text(encoding="utf-8").strip()
    if not text:
        raise ValueError(f"{path.name} is empty")
    return [RawSection(text=text, meta={"title": _make_title(text) or path.stem})]


def load_document(path: str | Path) -> list[RawSection]:
    p = Path(path)
    if not p.exists():
        raise FileNotFoundError(p)
    suffix = p.suffix.lower()
    if suffix in {".pptx"}:
        return load_pptx(p)
    if suffix in {".pdf"}:
        return load_pdf(p)
    if suffix in {".md", ".txt"}:
        return load_text(p)
    raise NotImplementedError(
        f"loader for {suffix} not implemented (MVP supports .pptx, text pdf, md and txt)"
    )
